from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette (matches dashboard dark theme) ----
BG = RGBColor(0x0A, 0x0E, 0x17)
CARD = RGBColor(0x14, 0x1B, 0x2A)
BORDER = RGBColor(0x2A, 0x33, 0x45)
TEXT_PRIMARY = RGBColor(0xF1, 0xF5, 0xF9)
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_RED = RGBColor(0xE6, 0x67, 0x67)
ACCENT_GREEN = RGBColor(0x0C, 0xA3, 0x0C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(slide, left, top, width, height, text, size=18, color=TEXT_PRIMARY,
             bold=False, align=PP_ALIGN.LEFT, font="Arial", line_spacing=1.15,
             anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_card(slide, left, top, width, height, fill=CARD, line_color=BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line_color
    card.line.width = Pt(1)
    card.shadow.inherit = False
    return card


def add_dot(slide, left, top, size, color):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    dot.shadow.inherit = False
    return dot


def header(slide, kicker, title, subtitle=None):
    add_dot(slide, Inches(0.5), Inches(0.42), Inches(0.14), ACCENT_GREEN)
    add_text(slide, Inches(0.75), Inches(0.32), Inches(6), Inches(0.35),
              kicker.upper(), size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12.3), Inches(0.9),
              title, size=32, color=TEXT_PRIMARY, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.5),
                  subtitle, size=15, color=TEXT_SECONDARY)


def footer(slide, page_num):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(4), Inches(0.3),
              "OilWatch AI — Caspian Hackathon 2026", size=10, color=TEXT_SECONDARY)
    add_text(slide, Inches(12.0), Inches(7.1), Inches(0.8), Inches(0.3),
              str(page_num), size=10, color=TEXT_SECONDARY, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1 — Title
# ============================================================
s = add_slide()
add_dot(s, Inches(5.9), Inches(2.15), Inches(0.22), ACCENT_GREEN)
add_text(s, Inches(0), Inches(2.5), SLIDE_W, Inches(1.1),
          "OilWatch AI", size=54, color=TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.5), SLIDE_W, Inches(0.6),
          "Обнаружение нефтяных разливов по спутниковым снимкам с помощью ИИ",
          size=20, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.3), SLIDE_W, Inches(0.5),
          "Caspian Sea Action Week — Хакатон 2026", size=14, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(6.3), SLIDE_W, Inches(0.5),
          "[Название команды]", size=18, color=TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(6.75), SLIDE_W, Inches(0.4),
          "[Участники команды]", size=13, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — Problem
# ============================================================
s = add_slide()
header(s, "Проблема", "Нефтяные разливы угрожают Каспийскому морю",
       "Каспий — замкнутый водоём: загрязнению некуда рассеиваться, течений, которые бы его вымывали, нет")
items = [
    ("Медленное обнаружение", "О разливах часто узнают с опозданием на часы или дни — случайно, от местных жителей, или когда ущерб уже виден."),
    ("Ручной мониторинг", "Патрули и визуальный осмотр не могут стабильно покрывать ~370 000 км² Каспия, тем более ночью или в непогоду."),
    ("Высокая экологическая цена", "Каспий окружён нефтедобывающей инфраструктурой; незамеченный на ранней стадии разлив превращается в экологическую и экономическую катастрофу."),
]
card_w = Inches(3.95)
gap = Inches(0.25)
left0 = Inches(0.5)
for i, (t, d) in enumerate(items):
    left = left0 + i * (card_w + gap)
    add_card(s, left, Inches(2.1), card_w, Inches(3.6))
    add_dot(s, left + Inches(0.3), Inches(2.45), Inches(0.16), ACCENT_RED)
    add_text(s, left + Inches(0.3), Inches(2.75), card_w - Inches(0.6), Inches(0.6),
              t, size=18, bold=True, color=TEXT_PRIMARY)
    add_text(s, left + Inches(0.3), Inches(3.35), card_w - Inches(0.6), Inches(2.1),
              d, size=13, color=TEXT_SECONDARY, line_spacing=1.3)
footer(s, 2)

# ============================================================
# SLIDE 3 — Goal
# ============================================================
s = add_slide()
header(s, "Цель", "Автоматически обнаруживать нефтяные разливы по реальным спутниковым данным")
add_card(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(3.6))
add_text(s, Inches(1.0), Inches(2.75), Inches(11.3), Inches(1.0),
          "Дать спасательным службам и регулирующим органам быстрый, недорогой и постоянно",
          size=20, color=TEXT_PRIMARY, line_spacing=1.4)
add_text(s, Inches(1.0), Inches(3.35), Inches(11.3), Inches(1.0),
          "доступный способ находить нефтяные разливы на Каспии — до того, как они распространятся.",
          size=20, color=TEXT_PRIMARY, line_spacing=1.4)
add_text(s, Inches(1.0), Inches(4.35), Inches(11.3), Inches(1.2),
          "Сегодня: анализ загруженных спутниковых снимков по запросу.\nЗавтра: непрерывный автоматический мониторинг с живого спутникового потока.",
          size=15, color=TEXT_SECONDARY, line_spacing=1.5)
footer(s, 3)

# ============================================================
# SLIDE 4 — Solution
# ============================================================
s = add_slide()
header(s, "Решение", "OilWatch AI — SAR-снимок на входе, обнаружение разлива на выходе")
steps = [
    ("1", "Загрузка", "SAR (радарный) спутниковый снимок загружается через дашборд."),
    ("2", "Анализ", "Модель сегментации U-Net классифицирует каждый пиксель: чистая вода или нефтяное пятно."),
    ("3", "Визуализация", "Дашборд показывает оригинал, цветовую маску, уверенность модели и площадь разлива."),
    ("4", "История", "Каждое сканирование сохраняется в историю для непрерывного мониторинга."),
]
card_w = Inches(2.95)
gap = Inches(0.15)
left0 = Inches(0.5)
for i, (num, t, d) in enumerate(steps):
    left = left0 + i * (card_w + gap)
    add_card(s, left, Inches(2.2), card_w, Inches(3.5))
    circ = slide_circ = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.25), Inches(2.5), Inches(0.55), Inches(0.55))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT_BLUE
    circ.line.fill.background()
    circ.shadow.inherit = False
    tf = circ.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = num
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    add_text(s, left + Inches(0.25), Inches(3.25), card_w - Inches(0.5), Inches(0.5),
              t, size=17, bold=True, color=TEXT_PRIMARY)
    add_text(s, left + Inches(0.25), Inches(3.75), card_w - Inches(0.5), Inches(1.8),
              d, size=12, color=TEXT_SECONDARY, line_spacing=1.3)
footer(s, 4)

# ============================================================
# SLIDE 5 — Target audience
# ============================================================
s = add_slide()
header(s, "Целевая аудитория", "Для кого создан продукт")
audiences = [
    ("Экологические ведомства", "Региональные и национальные органы, отвечающие за защиту Каспийского моря и реагирование на инциденты."),
    ("Нефтегазовые операторы", "Операторы морских платформ, которым нужно раннее предупреждение об утечках рядом со своей или соседней инфраструктурой."),
    ("Службы экстренного реагирования", "Спасатели, которым важно быстро узнать местоположение и масштаб разлива до начала ликвидации."),
    ("Исследователи и НПО", "Экологи и НПО, отслеживающие долгосрочные тренды загрязнения Каспия."),
]
card_w = Inches(5.9)
gap = Inches(0.5)
for i, (t, d) in enumerate(audiences):
    left = Inches(0.5) + (i % 2) * (card_w + gap)
    top = Inches(2.15) + (i // 2) * Inches(1.85)
    add_card(s, left, top, card_w, Inches(1.6))
    add_text(s, left + Inches(0.3), top + Inches(0.2), card_w - Inches(0.6), Inches(0.4),
              t, size=16, bold=True, color=TEXT_PRIMARY)
    add_text(s, left + Inches(0.3), top + Inches(0.65), card_w - Inches(0.6), Inches(0.85),
              d, size=12, color=TEXT_SECONDARY, line_spacing=1.25)
footer(s, 5)

# ============================================================
# SLIDE 6 — Key features
# ============================================================
s = add_slide()
header(s, "Ключевые функции", "Что умеет MVP уже сейчас")
features = [
    ("Загрузка и детекция", "Drag-and-drop загрузка SAR-снимка с мгновенным обнаружением разлива при помощи ИИ"),
    ("Маска сегментации", "Попиксельная визуализация того, где именно находится разлив, а не просто ответ да/нет"),
    ("Уверенность и площадь", "Численная оценка уверенности модели и площади разлива, а не просто метка"),
    ("История сканирований", "Каждое предсказание сохраняется и доступно для просмотра — журнал мониторинга"),
    ("Многоязычный интерфейс", "Дашборд доступен на английском, русском и казахском языках"),
    ("Запуск одной командой", "Весь стек (БД, backend, frontend, nginx) поднимается одной командой Docker"),
]
card_w = Inches(3.95)
gap = Inches(0.25)
for i, (t, d) in enumerate(features):
    left = Inches(0.5) + (i % 3) * (card_w + gap)
    top = Inches(2.15) + (i // 3) * Inches(1.95)
    add_card(s, left, top, card_w, Inches(1.7))
    add_dot(s, left + Inches(0.3), top + Inches(0.28), Inches(0.14), ACCENT_GREEN)
    add_text(s, left + Inches(0.55), top + Inches(0.2), card_w - Inches(0.85), Inches(0.4),
              t, size=14, bold=True, color=TEXT_PRIMARY)
    add_text(s, left + Inches(0.3), top + Inches(0.7), card_w - Inches(0.6), Inches(0.9),
              d, size=11, color=TEXT_SECONDARY, line_spacing=1.25)
footer(s, 6)

# ============================================================
# SLIDE 7 — Tech stack
# ============================================================
s = add_slide()
header(s, "Технологический стек", "Продакшн-архитектура, а не хакатонный костыль")
cols = [
    ("Backend", ["Python 3.13", "Django 5 + DRF", "TensorFlow / Keras", "PostgreSQL"]),
    ("Frontend", ["React + TypeScript", "Vite", "TailwindCSS", "EN / RU / KK i18n"]),
    ("Инфраструктура", ["Docker Compose", "nginx (единая точка входа)", "gunicorn", "Hugging Face Hub"]),
]
card_w = Inches(3.95)
gap = Inches(0.25)
for i, (title, rows) in enumerate(cols):
    left = Inches(0.5) + i * (card_w + gap)
    add_card(s, left, Inches(2.15), card_w, Inches(4.0))
    add_text(s, left + Inches(0.3), Inches(2.4), card_w - Inches(0.6), Inches(0.4),
              title, size=17, bold=True, color=ACCENT_BLUE)
    for j, row in enumerate(rows):
        yy = Inches(3.0) + j * Inches(0.62)
        add_dot(s, left + Inches(0.3), yy + Inches(0.07), Inches(0.09), TEXT_SECONDARY)
        add_text(s, left + Inches(0.55), yy, card_w - Inches(0.85), Inches(0.5),
                  row, size=13, color=TEXT_PRIMARY)
footer(s, 7)

# ============================================================
# SLIDE 8 — MVP description / what's real
# ============================================================
s = add_slide()
header(s, "Статус MVP", "Рабочий продукт, который можно независимо проверить")
left_w = Inches(6.0)
add_card(s, Inches(0.5), Inches(2.15), left_w, Inches(4.0))
add_text(s, Inches(0.8), Inches(2.4), left_w - Inches(0.6), Inches(0.4),
          "Что работает уже сейчас", size=16, bold=True, color=ACCENT_GREEN)
working = [
    "Реальный инференс через модель с Hugging Face (без имитации результатов)",
    "Проверено на реальных, опубликованных снимках разливов",
    "Полный деплой через Docker Compose: db + backend + frontend + nginx",
    "Демо-данные загружаются автоматически — дашборд не пустой при первом запуске",
    "Админ-панель для просмотра всех предсказаний",
]
for j, row in enumerate(working):
    yy = Inches(3.0) + j * Inches(0.62)
    add_dot(s, Inches(0.8), yy + Inches(0.07), Inches(0.09), ACCENT_GREEN)
    add_text(s, Inches(1.05), yy, left_w - Inches(0.85), Inches(0.55),
              row, size=12, color=TEXT_PRIMARY, line_spacing=1.2)

right_left = Inches(6.75)
right_w = Inches(6.05)
add_card(s, right_left, Inches(2.15), right_w, Inches(4.0))
add_text(s, right_left + Inches(0.3), Inches(2.4), right_w - Inches(0.6), Inches(0.4),
          "Честно задокументированные ограничения", size=16, bold=True, color=TEXT_SECONDARY)
limits = [
    "Нужны сырые SAR-снимки радара, а не обычные фото",
    "Пока только ручная загрузка — живой поток в планах",
    "Тренировочный датасет базовой модели официально не задокументирован",
]
for j, row in enumerate(limits):
    yy = Inches(3.0) + j * Inches(0.72)
    add_dot(s, right_left + Inches(0.3), yy + Inches(0.07), Inches(0.09), TEXT_SECONDARY)
    add_text(s, right_left + Inches(0.55), yy, right_w - Inches(0.85), Inches(0.65),
              row, size=12, color=TEXT_PRIMARY, line_spacing=1.2)
footer(s, 8)

# ============================================================
# SLIDE 9 — Growth potential
# ============================================================
s = add_slide()
header(s, "Потенциал роста", "Что будет после хакатона")
roadmap = [
    ("Ближайшие планы", "Автоматическая загрузка живого спутникового потока Sentinel-1 над Каспием — непрерывный мониторинг вместо ручной загрузки."),
    ("Среднесрочные планы", "Оповещения (SMS/email/Telegram) для служб реагирования при обнаружении разлива, с привязкой к карте."),
    ("Долгосрочные планы", "Дообучение модели на данных, специфичных для Каспия, и расширение на другие экологические угрозы (цветение водорослей, незаконный вылов рыбы)."),
]
card_w = Inches(3.95)
gap = Inches(0.25)
for i, (t, d) in enumerate(roadmap):
    left = Inches(0.5) + i * (card_w + gap)
    add_card(s, left, Inches(2.15), card_w, Inches(3.6))
    add_text(s, left + Inches(0.3), Inches(2.45), card_w - Inches(0.6), Inches(0.4),
              t, size=16, bold=True, color=ACCENT_BLUE)
    add_text(s, left + Inches(0.3), Inches(3.0), card_w - Inches(0.6), Inches(2.5),
              d, size=13, color=TEXT_SECONDARY, line_spacing=1.35)
add_text(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.6),
          "Обязуемся продолжать разработку до 31 декабря 2026 года — согласно условиям хакатона для победителей.",
          size=12, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
footer(s, 9)

# ============================================================
# SLIDE 10 — Live demo / closing
# ============================================================
s = add_slide()
add_dot(s, Inches(5.9), Inches(2.3), Inches(0.22), ACCENT_GREEN)
add_text(s, Inches(0), Inches(2.65), SLIDE_W, Inches(0.9),
          "Демонстрация", size=44, color=TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.55), SLIDE_W, Inches(0.5),
          "github.com/turkpenbayev/oil-watch", size=18, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.2), SLIDE_W, Inches(0.5),
          "docker compose up --build", size=16, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(5.3), SLIDE_W, Inches(0.5),
          "Вопросы?", size=20, color=TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

prs.save("/Users/bauka/Desktop/projects/oil-watch/docs/presentation/OilWatch_AI_Presentation.pptx")
print("Saved.")
